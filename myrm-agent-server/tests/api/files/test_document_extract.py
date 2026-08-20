"""Tests for document_extract API helpers and content_extraction parsers."""

from __future__ import annotations

import os
import tempfile
from pathlib import Path

import pytest

from app.api.files.document_extract import _validate_extension
from app.core.utils.errors import StandardHTTPException
from app.services.files.content_extraction import (
    _parse_document,
    extract_document_from_bytes,
)

_MIN_PNG = b"iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="


class TestValidateExtension:
    """Tests for _validate_extension."""

    def test_docx(self) -> None:
        assert _validate_extension("report.docx") == ".docx"

    def test_xlsx(self) -> None:
        assert _validate_extension("data.xlsx") == ".xlsx"

    def test_xls(self) -> None:
        assert _validate_extension("legacy.xls") == ".xls"

    def test_pptx(self) -> None:
        assert _validate_extension("slides.pptx") == ".pptx"

    def test_ppt(self) -> None:
        assert _validate_extension("old.ppt") == ".ppt"

    def test_unsupported_raises(self) -> None:
        with pytest.raises(StandardHTTPException):
            _validate_extension("image.png")

    def test_pdf_not_supported(self) -> None:
        with pytest.raises(StandardHTTPException):
            _validate_extension("doc.pdf")


class TestParseDocument:
    """Tests for _parse_document with real files."""

    @pytest.mark.asyncio
    async def test_parse_docx(self) -> None:
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            doc.add_paragraph("Extract test content")
            doc.save(f.name)
            tmp = f.name

        try:
            result = await _parse_document(Path(tmp), ".docx")
            assert "Extract test content" in result
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_extract_docx_embedded_image(self) -> None:
        import base64

        from docx import Document

        png_bytes = base64.b64decode(_MIN_PNG)
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as png:
                png.write(png_bytes)
                png_path = png.name
            try:
                doc.add_picture(png_path)
            finally:
                os.unlink(png_path)
            doc.save(f.name)
            tmp = f.name

        try:
            with open(tmp, "rb") as handle:
                content = handle.read()
            result = await extract_document_from_bytes(content, filename="pictured.docx")
            assert result.format == "docx"
            assert result.images, "docx embedded image should be extracted"
            assert result.images[0].mime_type in {"image/png", "image/jpeg"}
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_parse_xlsx(self) -> None:
        from openpyxl import Workbook

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = Workbook()
            ws = wb.active
            assert ws is not None
            ws.append(["Name", "Value"])
            ws.append(["Test", 42])
            wb.save(f.name)
            tmp = f.name

        try:
            result = await _parse_document(Path(tmp), ".xlsx")
            assert "Name" in result
            assert "Test" in result
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_parse_pptx(self) -> None:
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Extraction Test"
            slide.placeholders[1].text = "API endpoint works"
            prs.save(f.name)
            tmp = f.name

        try:
            result = await _parse_document(Path(tmp), ".pptx")
            assert "Extraction Test" in result
            assert "API endpoint works" in result
        finally:
            os.unlink(tmp)

    @pytest.mark.asyncio
    async def test_unsupported_extension_returns_empty(self) -> None:
        result = await _parse_document(Path("/tmp/fake.odt"), ".odt")
        assert result == ""
