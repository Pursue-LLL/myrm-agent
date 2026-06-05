"""Tests for agent_params/mention.py

Tests _build_mentioned_file_context, _inject_mentioned_files_into_query,
_parse_document, and _format_size with real file I/O.
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path

import pytest


class TestBuildMentionedFileContext:
    """Tests for _build_mentioned_file_context."""

    def _make_workspace(self, files: dict[str, str | bytes]) -> str:
        """Create a temporary workspace with given files."""
        tmpdir = tempfile.mkdtemp()
        for rel_path, content in files.items():
            abs_path = os.path.join(tmpdir, rel_path)
            os.makedirs(os.path.dirname(abs_path), exist_ok=True)
            mode = "wb" if isinstance(content, bytes) else "w"
            with open(abs_path, mode) as f:
                f.write(content)
        return tmpdir

    @pytest.mark.asyncio
    async def test_inline_small_text_file(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({"notes.txt": "Hello World"})
        result, warnings, tokens = await _build_mentioned_file_context(["notes.txt"], workspace)

        assert "<mentioned_files>" in result
        assert 'path="notes.txt"' in result
        assert 'type="text"' in result
        assert "Hello World" in result
        assert isinstance(warnings, list)
        assert isinstance(tokens, int)

    @pytest.mark.asyncio
    async def test_file_not_found(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({})
        result, warnings, tokens = await _build_mentioned_file_context(["missing.txt"], workspace)

        assert 'error="file not found"' in result

    @pytest.mark.asyncio
    async def test_directory_not_treated_as_file(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({"subdir/dummy.txt": "x"})
        result, warnings, tokens = await _build_mentioned_file_context(["subdir"], workspace)

        assert 'error="file not found"' in result

    @pytest.mark.asyncio
    async def test_path_outside_workspace(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({})
        result, warnings, tokens = await _build_mentioned_file_context(["../../etc/passwd"], workspace)

        assert 'error="path outside workspace"' in result

    @pytest.mark.asyncio
    async def test_binary_file_metadata_only(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({"image.png": b"\x89PNG\r\n\x1a\n" + b"\x00" * 100})
        result, warnings, tokens = await _build_mentioned_file_context(["image.png"], workspace)

        assert 'type="binary"' in result
        assert "use file_read_tool" in result

    @pytest.mark.asyncio
    async def test_pdf_binary_metadata_only(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({"doc.pdf": b"%PDF-1.4 fake content"})
        result, warnings, tokens = await _build_mentioned_file_context(["doc.pdf"], workspace)

        assert 'type="binary"' in result

    @pytest.mark.asyncio
    async def test_max_files_limit(self) -> None:
        from app.services.agent.params.mention import _MENTION_MAX_FILES, _build_mentioned_file_context

        files = {f"file_{i}.txt": f"content {i}" for i in range(15)}
        workspace = self._make_workspace(files)
        paths = [f"file_{i}.txt" for i in range(15)]
        result, warnings, tokens = await _build_mentioned_file_context(paths, workspace)

        assert result.count("mentioned_file") <= _MENTION_MAX_FILES * 2 + 2

    @pytest.mark.asyncio
    async def test_large_text_file_not_inlined(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        big_content = "x" * (150 * 1024)
        workspace = self._make_workspace({"big.txt": big_content})
        result, warnings, tokens = await _build_mentioned_file_context(["big.txt"], workspace)

        assert "Content too large to inline" in result

    @pytest.mark.asyncio
    async def test_multiple_files(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace(
            {
                "a.txt": "Alpha",
                "b.md": "Beta",
            }
        )
        result, warnings, tokens = await _build_mentioned_file_context(["a.txt", "b.md"], workspace)

        assert "Alpha" in result
        assert "Beta" in result

    @pytest.mark.asyncio
    async def test_empty_list_returns_empty(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_workspace({})
        result, warnings, tokens = await _build_mentioned_file_context([], workspace)

        assert result == ""

    @pytest.mark.asyncio
    async def test_docx_document_parsed(self) -> None:
        from docx import Document

        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        docx_path = os.path.join(workspace, "report.docx")
        doc = Document()
        doc.add_heading("Quarterly Report", level=1)
        doc.add_paragraph("Revenue increased 25%.")
        doc.save(docx_path)

        result, warnings, tokens = await _build_mentioned_file_context(["report.docx"], workspace)

        assert 'type="document"' in result
        assert "Quarterly Report" in result
        assert "Revenue increased 25%" in result

    @pytest.mark.asyncio
    async def test_xlsx_document_parsed(self) -> None:
        from openpyxl import Workbook

        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        xlsx_path = os.path.join(workspace, "data.xlsx")
        wb = Workbook()
        ws = wb.active
        assert ws is not None
        ws.append(["Name", "Score"])
        ws.append(["Alice", 95])
        wb.save(xlsx_path)

        result, warnings, tokens = await _build_mentioned_file_context(["data.xlsx"], workspace)

        assert 'type="document"' in result
        assert "Alice" in result

    @pytest.mark.asyncio
    async def test_pptx_document_parsed(self) -> None:
        from pptx import Presentation

        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        pptx_path = os.path.join(workspace, "slides.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Strategic Plan"
        slide.placeholders[1].text = "Key initiatives"
        prs.save(pptx_path)

        result, warnings, tokens = await _build_mentioned_file_context(["slides.pptx"], workspace)

        assert 'type="document"' in result
        assert "Strategic Plan" in result
        assert "Key initiatives" in result


class TestInjectMentionedFilesIntoQuery:
    """Tests for _inject_mentioned_files_into_query."""

    def test_string_query(self) -> None:
        from app.services.agent.params.mention import _inject_mentioned_files_into_query

        result = _inject_mentioned_files_into_query("Hello", "<ctx>files</ctx>")
        assert result == "Hello<ctx>files</ctx>"

    def test_empty_context_noop(self) -> None:
        from app.services.agent.params.mention import _inject_mentioned_files_into_query

        result = _inject_mentioned_files_into_query("Hello", "")
        assert result == "Hello"

    def test_multimodal_query_with_text(self) -> None:
        from app.services.agent.params.mention import _inject_mentioned_files_into_query

        query: list[dict[str, object]] = [{"type": "text", "text": "Analyze this"}]
        result = _inject_mentioned_files_into_query(query, "<ctx>data</ctx>")

        assert isinstance(result, list)
        assert "Analyze this<ctx>data</ctx>" in result[0]["text"]  # type: ignore[index]

    def test_multimodal_query_no_text_part(self) -> None:
        from app.services.agent.params.mention import _inject_mentioned_files_into_query

        query: list[dict[str, object]] = [{"type": "image_url", "url": "http://example.com/img.png"}]
        result = _inject_mentioned_files_into_query(query, "<ctx>data</ctx>")

        assert isinstance(result, list)
        assert len(result) == 2
        assert result[1]["type"] == "text"


class TestParseDocument:
    """Tests for _parse_document."""

    def test_docx_parse(self) -> None:
        from docx import Document

        from app.services.agent.params.mention import _parse_document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            doc = Document()
            doc.add_paragraph("Test paragraph")
            doc.save(f.name)
            result = _parse_document(f.name, ".docx")
            os.unlink(f.name)

        assert result is not None
        assert "Test paragraph" in result

    def test_xlsx_parse(self) -> None:
        from openpyxl import Workbook

        from app.services.agent.params.mention import _parse_document

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            wb = Workbook()
            ws = wb.active
            assert ws is not None
            ws.append(["Col1", "Col2"])
            wb.save(f.name)
            result = _parse_document(f.name, ".xlsx")
            os.unlink(f.name)

        assert result is not None
        assert "Col1" in result

    def test_pptx_parse(self) -> None:
        from pptx import Presentation

        from app.services.agent.params.mention import _parse_document

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Test Slide"
            prs.save(f.name)
            result = _parse_document(f.name, ".pptx")
            os.unlink(f.name)

        assert result is not None
        assert "Test Slide" in result

    def test_unsupported_extension(self) -> None:
        from app.services.agent.params.mention import _parse_document

        result = _parse_document("/tmp/fake.odt", ".odt")
        assert result is None


class TestFormatSize:
    """Tests for _format_size."""

    def test_bytes(self) -> None:
        from app.services.agent.params.mention import _format_size

        assert _format_size(500) == "500B"

    def test_kilobytes(self) -> None:
        from app.services.agent.params.mention import _format_size

        assert _format_size(2048) == "2.0KB"

    def test_megabytes(self) -> None:
        from app.services.agent.params.mention import _format_size

        assert _format_size(5 * 1024 * 1024) == "5.0MB"


class TestRichContextReferences:
    """Tests for @staged, @folder:, @url: references."""

    def _make_git_workspace(self, files: dict[str, str]) -> str:
        """Create a temporary git workspace with staged files."""
        import subprocess

        tmpdir = tempfile.mkdtemp()
        subprocess.run(["git", "init"], cwd=tmpdir, check=True, capture_output=True)
        subprocess.run(["git", "config", "user.email", "test@test.com"], cwd=tmpdir, check=True, capture_output=True)
        subprocess.run(["git", "config", "user.name", "Test User"], cwd=tmpdir, check=True, capture_output=True)

        for rel_path, content in files.items():
            abs_path = os.path.join(tmpdir, rel_path)
            os.makedirs(os.path.dirname(abs_path), exist_ok=True)
            with open(abs_path, "w") as f:
                f.write(content)

        return tmpdir

    @pytest.mark.asyncio
    async def test_staged_reference_no_changes(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_git_workspace({"test.txt": "hello"})
        result, warnings, tokens = await _build_mentioned_file_context(["@staged"], workspace)

        assert "<mentioned_file" in result
        assert 'path="@staged"' in result
        assert 'type="git-diff"' in result
        assert "No staged changes" in result

    @pytest.mark.asyncio
    async def test_staged_reference_with_changes(self) -> None:
        import subprocess

        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_git_workspace({"test.txt": "original"})
        subprocess.run(["git", "add", "test.txt"], cwd=workspace, check=True, capture_output=True)
        subprocess.run(["git", "commit", "-m", "initial"], cwd=workspace, check=True, capture_output=True)

        with open(os.path.join(workspace, "test.txt"), "w") as f:
            f.write("modified")
        subprocess.run(["git", "add", "test.txt"], cwd=workspace, check=True, capture_output=True)

        result, warnings, tokens = await _build_mentioned_file_context(["@staged"], workspace)

        assert 'path="@staged"' in result
        assert 'type="git-diff"' in result
        assert "diff --git" in result or "+modified" in result

    @pytest.mark.asyncio
    async def test_folder_reference_empty(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        result, warnings, tokens = await _build_mentioned_file_context(["@folder:"], workspace)

        assert "<mentioned_file" in result
        assert 'path="@folder:"' in result

    @pytest.mark.asyncio
    async def test_folder_reference_with_files(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        os.makedirs(os.path.join(workspace, "src"))
        os.makedirs(os.path.join(workspace, "src", "utils"))
        Path(os.path.join(workspace, "src", "main.py")).write_text("print('hello')")
        Path(os.path.join(workspace, "src", "utils", "helper.py")).write_text("def helper(): pass")

        result, warnings, tokens = await _build_mentioned_file_context(["@folder:src"], workspace)

        assert 'path="@folder:src"' in result
        assert 'type="folder-tree"' in result
        assert "main.py" in result
        assert "utils/" in result
        assert "helper.py" in result

    @pytest.mark.asyncio
    async def test_folder_reference_outside_workspace(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        result, warnings, tokens = await _build_mentioned_file_context(["@folder:../../etc"], workspace)

        assert 'error="path outside workspace"' in result

    @pytest.mark.asyncio
    async def test_folder_reference_filters_excluded_dirs(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        os.makedirs(os.path.join(workspace, ".git"))
        os.makedirs(os.path.join(workspace, "node_modules"))
        os.makedirs(os.path.join(workspace, "src"))
        Path(os.path.join(workspace, ".git", "config")).write_text("git config")
        Path(os.path.join(workspace, "node_modules", "package.json")).write_text("{}")
        Path(os.path.join(workspace, "src", "main.py")).write_text("code")

        result, warnings, tokens = await _build_mentioned_file_context(["@folder:"], workspace)

        assert "main.py" in result
        assert ".git" not in result
        assert "node_modules" not in result

    @pytest.mark.asyncio
    async def test_url_reference_invalid_url(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = tempfile.mkdtemp()
        result, warnings, tokens = await _build_mentioned_file_context(["@url:not-a-valid-url"], workspace)

        assert 'path="@url:not-a-valid-url"' in result
        assert 'error="failed to fetch URL"' in result

    @pytest.mark.asyncio
    async def test_mixed_references(self) -> None:
        from app.services.agent.params.mention import _build_mentioned_file_context

        workspace = self._make_git_workspace({"test.txt": "hello", "data.md": "# Data"})

        result, warnings, tokens = await _build_mentioned_file_context(["@staged", "test.txt", "@folder:"], workspace)

        assert 'path="@staged"' in result
        assert 'path="test.txt"' in result
        assert 'path="@folder:"' in result
        assert "hello" in result
